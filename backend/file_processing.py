import fitz  # PyMuPDF for PDF extraction
import pptx

def extract_text_from_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
    return text
